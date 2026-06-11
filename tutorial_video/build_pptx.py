#!/usr/bin/env python3
"""Build a standalone 16:9 PPTX from the rendered deck scenes.
One full-bleed slide image per scene; the scene caption becomes the speaker note."""
import re
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches

TV = Path("/Users/huangshifeng/Desktop/Research/Oncology_Outcomes/stage_III_colon_edr/tutorial_video")
FRAMES = TV / "pptx_frames"
OUT = TV / "output" / "tutorial_slides.pptx"


def captions():
    html = (TV / "deck" / "index.html").read_text(encoding="utf-8")
    body = re.search(r"const CAPS\s*=\s*\[(.*?)\];", html, re.S).group(1)
    caps = re.findall(r'"((?:[^"\\]|\\.)*)"', body)
    return caps


def main():
    caps = captions()
    pngs = sorted(FRAMES.glob("scene_*.png"))
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]
    for i, png in enumerate(pngs):
        slide = prs.slides.add_slide(blank)
        slide.shapes.add_picture(str(png), 0, 0, width=prs.slide_width,
                                 height=prs.slide_height)
        note = caps[i] if i < len(caps) else ""
        if note:
            slide.notes_slide.notes_text_frame.text = note
    OUT.parent.mkdir(exist_ok=True)
    prs.save(str(OUT))
    print(f"wrote {OUT} with {len(pngs)} slides")


if __name__ == "__main__":
    main()
